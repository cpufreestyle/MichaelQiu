import argparse
import os
import sys
import re
from docx import Document
from openpyxl import load_workbook
from PyPDF2 import PdfReader, PdfWriter
from pptx import Presentation
import warnings
import contextlib
import csv

# 屏蔽PyPDF2及相关库的所有警告输出
warnings.filterwarnings("ignore", category=UserWarning)
warnings.filterwarnings("ignore", category=DeprecationWarning)
try:
    from PyPDF2.errors import PdfReadWarning
    warnings.filterwarnings("ignore", category=PdfReadWarning)
except Exception:
    pass

# 支持的文件类型
FILE_TYPES = ['.docx', '.xlsx', '.pdf', '.ppt', '.pptx']

# 获取所有磁盘根目录（Windows）
def get_all_drives():
    import string
    from ctypes import windll
    drives = []
    DRIVE_REMOVABLE = 2
    DRIVE_CDROM = 5
    DRIVE_FIXED = 3
    bitmask = windll.kernel32.GetLogicalDrives()
    for letter in string.ascii_uppercase:
        if bitmask & 1:
            drive = f'{letter}:\\'
            # 排除外置存储（如手机、U盘、SD卡等）
            try:
                drive_type = windll.kernel32.GetDriveTypeW(drive)
                # 2: removable, 5: cdrom
                if drive_type in [DRIVE_REMOVABLE, DRIVE_CDROM]:
                    bitmask >>= 1
                    continue
            except Exception:
                pass
            drives.append(drive)
        bitmask >>= 1
    return drives

# 处理docx文件
def process_docx(file_path, success_files, changed_files):
    try:
        doc = Document(file_path)
        changed = False
        if doc.paragraphs:
            new_text = doc.paragraphs[0].text.replace('内部', '')
            if new_text != doc.paragraphs[0].text:
                doc.paragraphs[0].text = new_text
                changed = True
        doc.save(file_path)
        success_files.append(file_path)
        if changed:
            changed_files.append(file_path)
    except Exception:
        pass

# 处理xlsx文件
def process_xlsx(file_path, success_files, changed_files):
    try:
        wb = load_workbook(file_path)
        ws = wb.active
        first_row = next(ws.iter_rows(min_row=1, max_row=1))
        changed = False
        for cell in first_row:
            if cell.value and isinstance(cell.value, str):
                new_val = cell.value.replace('内部', '')
                if new_val != cell.value:
                    cell.value = new_val
                    changed = True
        wb.save(file_path)
        success_files.append(file_path)
        if changed:
            changed_files.append(file_path)
    except Exception:
        pass

# 处理pdf文件
def process_pdf(file_path, success_files, changed_files):
    try:
        with open(os.devnull, 'w') as devnull, contextlib.redirect_stderr(devnull):
            reader = PdfReader(file_path)
            writer = PdfWriter()
            changed = False
            if reader.pages:
                first_page = reader.pages[0]
                text = first_page.extract_text()
                if text:
                    first_page_content = first_page.extract_text().split('\n')
                    if first_page_content:
                        new_line = first_page_content[0].replace('内部', '')
                        if new_line != first_page_content[0]:
                            changed = True
                        first_page_content[0] = new_line
                writer.add_page(first_page)
                for page in reader.pages[1:]:
                    writer.add_page(page)
                with open(file_path, 'wb') as f:
                    writer.write(f)
            success_files.append(file_path)
            if changed:
                changed_files.append(file_path)
    except Exception:
        pass

# 处理ppt/pptx文件
def process_ppt(file_path, success_files, changed_files):
    try:
        prs = Presentation(file_path)
        changed = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines = shape.text.split('\n')
                    if lines:
                        new_line = lines[0].replace('内部', '')
                        if new_line != lines[0]:
                            changed = True
                        lines[0] = new_line
                        shape.text = '\n'.join(lines)
        prs.save(file_path)
        success_files.append(file_path)
        if changed:
            changed_files.append(file_path)
    except Exception:
        pass

# 递归扫描并处理文件
def scan_and_process(root_dir, success_files, changed_files, file_counter):
    for dirpath, _, filenames in os.walk(root_dir):
        # 跳过OneDrive等网盘目录和指定目录
        skip_keywords = ["onedrive", "cloud", "dropbox", "baidunetdisk", "坚果云", "sync", "googledrive"]
        skip_paths = [r"c:\\users\\michael\\crossdevice\\michael  s ultra\\"]
        dirpath_lower = dirpath.lower()
        if any(x in dirpath_lower for x in skip_keywords) or any(dirpath_lower.startswith(p) for p in skip_paths):
            continue
        for filename in filenames:
            ext = os.path.splitext(filename)[1].lower()
            if ext in FILE_TYPES:
                file_path = os.path.join(dirpath, filename)
                if ext == '.docx':
                    process_docx(file_path, success_files, changed_files)
                elif ext == '.xlsx':
                    process_xlsx(file_path, success_files, changed_files)
                elif ext == '.pdf':
                    process_pdf(file_path, success_files, changed_files)
                elif ext in ['.ppt', '.pptx']:
                    process_ppt(file_path, success_files, changed_files)
                file_counter[0] += 1
                if file_counter[0] % 100 == 0:
                    print(f"已处理文件数: {file_counter[0]}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="递归扫描并处理文档")
    parser.add_argument('--dir', type=str, default=None, help='只扫描指定目录（如C:/Users/michael）')
    args = parser.parse_args()
    drives = get_all_drives() if args.dir is None else [args.dir]
    success_files = []
    changed_files = []
    file_counter = [0]
    for drive in drives:
        print(f"扫描: {drive}")
        scan_and_process(drive, success_files, changed_files, file_counter)
    print("处理完成！")
    print("成功修改的文件列表：")
    for f in success_files:
        print(f)
    print("\n被删除‘内部’的文件列表已导出为changed_files.csv")
    with open("changed_files.csv", "w", newline='', encoding="utf-8") as csvfile:
        writer = csv.writer(csvfile)
        writer.writerow(["file_path"])
        for f in changed_files:
            writer.writerow([f])
